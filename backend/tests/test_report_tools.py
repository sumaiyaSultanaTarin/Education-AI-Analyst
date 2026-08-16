from docx import Document
from pptx import Presentation

from tools.report_tools import build_report_docx, build_report_pptx, summarize_social_intelligence


def test_build_report_docx_writes_a_readable_file(tmp_path):
    output_path = tmp_path / "session" / "report.docx"

    build_report_docx(
        str(output_path),
        goal="Summarize Spring 2025 performance",
        documents=[{"filename": "results.xlsx", "type": "xlsx"}],
        data_analysis={"doc-1": {"Term1": {"score": {"mean": 70, "min": 30, "max": 100, "pass_rate": 80, "count": 5}}}},
        rag_citations=[{"filename": "results.xlsx", "text": "excerpt text"}],
    )

    assert output_path.exists()

    doc = Document(str(output_path))
    full_text = "\n".join(p.text for p in doc.paragraphs)
    assert "Summarize Spring 2025 performance" in full_text
    assert "results.xlsx" in full_text
    assert "excerpt text" in full_text


def test_build_report_docx_omits_empty_sections(tmp_path):
    output_path = tmp_path / "report.docx"

    build_report_docx(
        str(output_path), goal="g", documents=[], data_analysis={}, rag_citations=[]
    )

    doc = Document(str(output_path))
    headings = [p.text for p in doc.paragraphs if p.style.name.startswith("Heading")]
    assert "Data Analysis" not in headings
    assert "Supporting Excerpts" not in headings


def test_build_report_pptx_writes_a_readable_file(tmp_path):
    output_path = tmp_path / "session" / "report.pptx"

    build_report_pptx(
        str(output_path),
        goal="Summarize Spring 2025 performance",
        documents=[{"filename": "results.xlsx", "type": "xlsx"}],
        data_analysis={"doc-1": {"Term1": {"score": {"mean": 70, "min": 30, "max": 100, "pass_rate": 80, "count": 5}}}},
        rag_citations=[{"filename": "results.xlsx", "text": "excerpt text"}],
        social_summary={"post_count": 2, "comment_count": 3, "positive": 2, "neutral": 1, "negative": 0},
    )

    assert output_path.exists()

    prs = Presentation(str(output_path))
    all_text = "\n".join(
        shape.text_frame.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame
    )
    assert "Summarize Spring 2025 performance" in all_text
    assert "results.xlsx" in all_text
    assert "excerpt text" in all_text
    assert "Positive: 2" in all_text
    # title + documents + data analysis + social + citations
    assert len(prs.slides) == 5


def test_summarize_social_intelligence_counts_sentiment_across_documents():
    social_outputs = {
        "doc-1": {"posts": [{
            "fb_post_id": "p1", "content": "x", "posted_at": "t",
            "comments": [
                {"sentiment": {"label": "positive"}},
                {"sentiment": {"label": "negative"}},
            ],
        }]},
        "doc-2": {"posts": [{
            "fb_post_id": "p2", "content": "y", "posted_at": "t",
            "comments": [{"sentiment": {"label": "positive"}}],
        }]},
    }

    summary = summarize_social_intelligence(social_outputs)

    assert summary == {
        "post_count": 2, "comment_count": 3, "positive": 2, "neutral": 0, "negative": 1,
    }


def test_summarize_social_intelligence_returns_none_when_empty():
    assert summarize_social_intelligence({}) is None
