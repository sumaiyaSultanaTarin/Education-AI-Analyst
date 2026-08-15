from pptx import Presentation

from tools.pptx_tools import extract_pptx_content


def test_extract_pptx_content(tmp_path):
    # No sample .pptx ships in data/sample_docs/ (Phase 1 pack covers xlsx/pdf/docx/
    # images only), so build a minimal deck here instead.
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Term Overview"
    slide.placeholders[1].text_frame.text = "Enrollment up 6% this term."
    path = tmp_path / "sample.pptx"
    prs.save(path)

    result = extract_pptx_content(str(path))

    assert len(result["slides"]) == 1
    slide_text = result["slides"][0]["text"]
    assert any("Term Overview" in t for t in slide_text)
    assert any("Enrollment up 6%" in t for t in slide_text)
