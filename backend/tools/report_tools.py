"""DOCX/PPTX report assembly, used by the Report Generator Agent."""

from pathlib import Path

from docx import Document
from pptx import Presentation


def summarize_social_intelligence(social_outputs: dict) -> dict | None:
    """Roll the Social Intelligence Agent's per-document output up into
    post/comment/sentiment counts for the report.

    `social_outputs` is keyed by document_id -> {"posts": [...]} (see
    agents/social_intel_agent.py). Returns None if there's nothing to
    summarize, so callers can skip the report section entirely.
    """
    post_count = 0
    sentiment_counts = {"positive": 0, "neutral": 0, "negative": 0}

    for parsed in social_outputs.values():
        posts = parsed.get("posts", [])
        post_count += len(posts)
        for post in posts:
            for comment in post.get("comments", []):
                label = comment.get("sentiment", {}).get("label")
                if label in sentiment_counts:
                    sentiment_counts[label] += 1

    comment_count = sum(sentiment_counts.values())
    if post_count == 0 and comment_count == 0:
        return None

    return {"post_count": post_count, "comment_count": comment_count, **sentiment_counts}


def _data_analysis_lines(data_analysis: dict) -> list[str]:
    lines = []
    for summary in data_analysis.values():
        for sheet_name, stats in summary.items():
            for column, values in stats.items():
                lines.append(
                    f"{sheet_name} / {column}: mean={values['mean']}, min={values['min']}, "
                    f"max={values['max']}, pass_rate={values['pass_rate']}%, n={values['count']}"
                )
    return lines


def build_report_docx(
    output_path: str,
    goal: str,
    documents: list[dict],
    data_analysis: dict,
    rag_citations: list[dict],
    social_summary: dict | None = None,
    web_context: list[dict] | None = None,
) -> str:
    """Compile a formatted DOCX report with citations back to source documents.

    `data_analysis` is keyed by document_id -> {sheet_name: column stats}
    (Data Analyst Agent's output shape) — must NOT include a "web_context"
    key (callers split that out separately, see agents/report_generator_agent.py),
    since it's a list, not a {sheet_name: stats} dict, and every value here
    gets iterated the same way. `rag_citations` is the Knowledge/RAG Agent's
    query hits; `social_summary` is summarize_social_intelligence()'s output;
    `web_context` is tools/web_search_tools.search_web()'s output. Returns
    the path written to.
    """
    doc = Document()
    doc.add_heading("Education AI Analyst — Report", level=0)
    doc.add_paragraph(f"Goal: {goal}")

    doc.add_heading("Source Documents", level=1)
    for document in documents:
        doc.add_paragraph(f"{document['filename']} ({document['type']})", style="List Bullet")

    if data_analysis:
        doc.add_heading("Data Analysis", level=1)
        for summary in data_analysis.values():
            for sheet_name, stats in summary.items():
                doc.add_heading(sheet_name, level=2)
                for column, values in stats.items():
                    doc.add_paragraph(
                        f"{column}: mean={values['mean']}, min={values['min']}, "
                        f"max={values['max']}, pass_rate={values['pass_rate']}%, "
                        f"n={values['count']}"
                    )

    if social_summary:
        doc.add_heading("Social Intelligence (Facebook)", level=1)
        doc.add_paragraph(
            f"{social_summary['post_count']} post(s), {social_summary['comment_count']} comment(s) — "
            f"{social_summary['positive']} positive, {social_summary['neutral']} neutral, "
            f"{social_summary['negative']} negative."
        )

    if rag_citations:
        doc.add_heading("Supporting Excerpts", level=1)
        for citation in rag_citations:
            doc.add_paragraph(f"[{citation['filename']}] {citation['text'][:300]}")

    if web_context:
        doc.add_heading("External Benchmark Context", level=1)
        for result in web_context:
            doc.add_paragraph(f"{result['title']} — {result['snippet']}", style="List Bullet")
            doc.add_paragraph(result["url"])

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    doc.save(output_path)
    return output_path


def build_report_pptx(
    output_path: str,
    goal: str,
    documents: list[dict],
    data_analysis: dict,
    rag_citations: list[dict],
    social_summary: dict | None = None,
    web_context: list[dict] | None = None,
) -> str:
    """Compile the same report content as build_report_docx() as a slide deck.

    Same parameters/shapes as build_report_docx(). Returns the path written to.
    """
    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]

    title_slide = prs.slides.add_slide(title_layout)
    title_slide.shapes.title.text = "Education AI Analyst — Report"
    title_slide.placeholders[1].text = f"Goal: {goal}"

    docs_slide = prs.slides.add_slide(content_layout)
    docs_slide.shapes.title.text = "Source Documents"
    _fill_bullets(
        docs_slide.placeholders[1], [f"{d['filename']} ({d['type']})" for d in documents]
    )

    if data_analysis:
        analysis_slide = prs.slides.add_slide(content_layout)
        analysis_slide.shapes.title.text = "Data Analysis"
        _fill_bullets(analysis_slide.placeholders[1], _data_analysis_lines(data_analysis))

    if social_summary:
        social_slide = prs.slides.add_slide(content_layout)
        social_slide.shapes.title.text = "Social Intelligence (Facebook)"
        _fill_bullets(social_slide.placeholders[1], [
            f"{social_summary['post_count']} post(s), {social_summary['comment_count']} comment(s)",
            f"Positive: {social_summary['positive']}",
            f"Neutral: {social_summary['neutral']}",
            f"Negative: {social_summary['negative']}",
        ])

    if rag_citations:
        citations_slide = prs.slides.add_slide(content_layout)
        citations_slide.shapes.title.text = "Supporting Excerpts"
        _fill_bullets(
            citations_slide.placeholders[1],
            [f"[{c['filename']}] {c['text'][:200]}" for c in rag_citations],
        )

    if web_context:
        context_slide = prs.slides.add_slide(content_layout)
        context_slide.shapes.title.text = "External Benchmark Context"
        _fill_bullets(
            context_slide.placeholders[1],
            [f"{r['title']} — {r['snippet'][:150]}" for r in web_context],
        )

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


def _fill_bullets(placeholder, lines: list[str]) -> None:
    if not lines:
        placeholder.text_frame.text = "(none)"
        return

    placeholder.text_frame.text = lines[0]
    for line in lines[1:]:
        placeholder.text_frame.add_paragraph().text = line
