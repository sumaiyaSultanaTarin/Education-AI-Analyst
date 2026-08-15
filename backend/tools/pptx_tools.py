"""PPTX text extraction, used by the Document Ingestion Agent."""

from pptx import Presentation


def extract_pptx_content(path: str) -> dict:
    """Extract text from every shape on every slide.

    Returns {"slides": [{"slide_number": int, "text": list[str]}]}.
    """
    prs = Presentation(path)

    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = [
            shape.text_frame.text
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        slides.append({"slide_number": i, "text": texts})

    return {"slides": slides}
